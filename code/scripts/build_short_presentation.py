#!/usr/bin/env python3
from __future__ import annotations

import json
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
PAPER_FIGURES = ROOT / "paper" / "figures"
METRICS_DIR = ROOT / "results" / "metrics"
METADATA_DIR = ROOT / "metadata" / "lidc_index"
OUTPUT_DIR = ROOT / "presentation"
OUTPUT_PPTX = OUTPUT_DIR / "ldct_hybrid_egldm_3min_presentation.pptx"
OUTPUT_SCRIPT = OUTPUT_DIR / "ldct_hybrid_egldm_3min_script.txt"


TITLE_COLOR = RGBColor(18, 31, 53)
BODY_COLOR = RGBColor(42, 52, 65)
ACCENT = RGBColor(18, 114, 136)
ACCENT_SOFT = RGBColor(226, 241, 244)
GOLD = RGBColor(205, 164, 82)
MUTED = RGBColor(110, 118, 130)
LIGHT = RGBColor(246, 249, 252)
WHITE = RGBColor(255, 255, 255)
GREEN_SOFT = RGBColor(230, 244, 236)
RED_SOFT = RGBColor(248, 235, 233)


def load_json(path: Path) -> dict:
    return json.loads(path.read_text())


def rgb_tuple(color: RGBColor) -> tuple[int, int, int]:
    return color[0], color[1], color[2]


def add_rect(slide, x, y, w, h, fill, line=None, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def add_textbox(
    slide,
    text: str,
    x,
    y,
    w,
    h,
    *,
    font_size=20,
    color=BODY_COLOR,
    bold=False,
    name="Aptos",
    align=PP_ALIGN.LEFT,
    fill=None,
    margin=0.08,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.color.rgb = fill
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = name
    font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = color
    return box


def add_bullets(slide, items, x, y, w, h, *, font_size=20, color=BODY_COLOR, bullet_color=ACCENT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.alignment = PP_ALIGN.LEFT
        font = p.font
        font.name = "Aptos"
        font.size = Pt(font_size)
        font.color.rgb = color
        font.bold = False
    return box


def add_title(slide, title: str, subtitle: str | None = None):
    add_textbox(slide, title, Inches(0.72), Inches(0.35), Inches(11.0), Inches(0.8), font_size=28, color=TITLE_COLOR, bold=True)
    if subtitle:
        add_textbox(slide, subtitle, Inches(0.74), Inches(1.0), Inches(11.2), Inches(0.45), font_size=13, color=MUTED)
    bar = add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.16), ACCENT, ACCENT, radius=False)
    bar.line.fill.background()


def add_slide_number(slide, num: int):
    add_textbox(slide, str(num), Inches(12.7), Inches(7.03), Inches(0.3), Inches(0.2), font_size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def add_metric_card(slide, title: str, value: str, x, y, w, h, fill, value_color=TITLE_COLOR):
    add_rect(slide, x, y, w, h, fill, line=fill)
    add_textbox(slide, title, x + Inches(0.08), y + Inches(0.08), w - Inches(0.16), Inches(0.22), font_size=11, color=MUTED, bold=True)
    add_textbox(slide, value, x + Inches(0.08), y + Inches(0.34), w - Inches(0.16), h - Inches(0.34), font_size=22, color=value_color, bold=True)


def add_image_contain(slide, image_path: Path, x, y, w, h):
    with Image.open(image_path) as im:
        img_w, img_h = im.size
    box_w = float(w)
    box_h = float(h)
    img_ratio = img_w / img_h
    box_ratio = box_w / box_h
    if img_ratio > box_ratio:
        width = box_w
        height = width / img_ratio
        left = float(x)
        top = float(y) + (box_h - height) / 2
    else:
        height = box_h
        width = height * img_ratio
        top = float(y)
        left = float(x) + (box_w - width) / 2
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def add_table_like_results(slide, rows, x, y, w, h):
    cols = [0.39, 0.19, 0.19, 0.19]
    row_h = h / (len(rows) + 1)
    header_fill = ACCENT
    headers = ["Method", "PSNR", "SSIM", "RMSE"]
    cur_x = x
    for idx, head in enumerate(headers):
        cell_w = w * cols[idx]
        add_rect(slide, cur_x, y, cell_w, row_h, header_fill, line=header_fill)
        add_textbox(slide, head, cur_x, y + Inches(0.02), cell_w, row_h, font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        cur_x += cell_w

    for r_idx, row in enumerate(rows):
        row_y = y + row_h * (r_idx + 1)
        highlight = row["highlight"]
        fill = ACCENT_SOFT if highlight else WHITE
        line = RGBColor(218, 224, 232)
        cur_x = x
        vals = [row["name"], row["psnr"], row["ssim"], row["rmse"]]
        for c_idx, val in enumerate(vals):
            cell_w = w * cols[c_idx]
            add_rect(slide, cur_x, row_y, cell_w, row_h, fill, line=line)
            align = PP_ALIGN.LEFT if c_idx == 0 else PP_ALIGN.CENTER
            text_x = cur_x + (Inches(0.05) if c_idx == 0 else Inches(0))
            text_w = cell_w - (Inches(0.08) if c_idx == 0 else Inches(0))
            add_textbox(
                slide,
                val,
                text_x,
                row_y + Inches(0.01),
                text_w,
                row_h,
                font_size=12,
                color=TITLE_COLOR if highlight else BODY_COLOR,
                bold=highlight,
                align=align,
            )
            cur_x += cell_w


def build_presentation():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    hybrid = load_json(METRICS_DIR / "hybrid_test_best_metrics.json")["metrics"]
    redcnn = load_json(METRICS_DIR / "redcnn_test_metrics.json")["metrics"]
    edge = load_json(METRICS_DIR / "edge_guided_ldm_test_metrics.json")["metrics"]
    no_edge = load_json(METRICS_DIR / "no_edge_ldm_test_metrics.json")["metrics"]
    split = load_json(METADATA_DIR / "split_manifest.json")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    add_title(
        slide,
        "Hybrid Edge-Guided Diffusion Refinement for Leakage-Free LDCT Denoising",
        "Zhexi Feng, Bingrui Zhang | UC San Diego | 3-minute conference presentation",
    )
    add_textbox(
        slide,
        "Main message: diffusion is most effective here as a conservative anatomical refiner around a strong deterministic RED-CNN anchor.",
        Inches(0.78),
        Inches(1.45),
        Inches(11.8),
        Inches(0.72),
        font_size=19,
        color=TITLE_COLOR,
        fill=ACCENT_SOFT,
        bold=True,
    )
    add_metric_card(slide, "Patients", f"{split['scan_stats']['patients']:,}", Inches(0.82), Inches(2.45), Inches(2.55), Inches(1.2), GREEN_SOFT)
    add_metric_card(slide, "Valid CT slices", f"{split['scan_stats']['valid_slices']:,}", Inches(3.6), Inches(2.45), Inches(2.9), Inches(1.2), GREEN_SOFT)
    add_metric_card(slide, "Hybrid test result", "35.12 / 0.904 / 0.0356", Inches(6.73), Inches(2.45), Inches(4.6), Inches(1.2), RED_SOFT, value_color=ACCENT)
    add_bullets(
        slide,
        [
            "Leakage-free patient-level LIDC-IDRI protocol on raw DICOM.",
            "Hybrid EG-LDM slightly exceeds the RED-CNN anchor on the held-out test split.",
            "Standalone diffusion is weaker; hybrid refinement is the key result.",
        ],
        Inches(0.95),
        Inches(4.15),
        Inches(7.2),
        Inches(2.2),
        font_size=18,
    )
    add_image_contain(slide, PAPER_FIGURES / "hybrid_roi.png", Inches(8.55), Inches(4.0), Inches(4.25), Inches(2.55))
    add_slide_number(slide, 1)

    # Slide 2
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Why This Benchmark Matters")
    add_textbox(slide, "Problem", Inches(0.82), Inches(1.12), Inches(2.0), Inches(0.3), font_size=18, color=ACCENT, bold=True)
    add_bullets(
        slide,
        [
            "Low-dose CT denoising must reduce noise without inventing anatomy.",
            "Slice-level random splits in volumetric CT can leak near-duplicate anatomy.",
            "LIDC-IDRI has real anatomy but not paired standard-dose / low-dose scans.",
        ],
        Inches(0.86),
        Inches(1.42),
        Inches(5.55),
        Inches(2.35),
        font_size=18,
    )
    add_textbox(slide, "Protocol", Inches(0.82), Inches(4.05), Inches(2.0), Inches(0.3), font_size=18, color=ACCENT, bold=True)
    add_bullets(
        slide,
        [
            "Scan 382,009 candidate DICOM files and retain 243,957 valid CT slices from 1,010 patients.",
            "Persist patient-level train / val / test splits: 808 / 101 / 101 patients.",
            "Construct paired supervision by adding signal-dependent Gaussian corruption to clean preprocessed CT slices.",
        ],
        Inches(0.86),
        Inches(4.35),
        Inches(6.0),
        Inches(2.15),
        font_size=17,
    )
    add_metric_card(slide, "Train / Val / Test patients", "808 / 101 / 101", Inches(7.3), Inches(1.35), Inches(2.4), Inches(1.0), ACCENT_SOFT)
    add_metric_card(slide, "Series", f"{split['scan_stats']['series']:,}", Inches(9.95), Inches(1.35), Inches(1.6), Inches(1.0), ACCENT_SOFT)
    add_metric_card(slide, "Valid slices", f"{split['scan_stats']['valid_slices']:,}", Inches(11.8), Inches(1.35), Inches(1.0), Inches(1.0), ACCENT_SOFT)
    add_rect(slide, Inches(7.28), Inches(2.75), Inches(5.1), Inches(2.75), LIGHT, line=RGBColor(225, 229, 235))
    add_textbox(slide, "Preprocessing pipeline", Inches(7.52), Inches(2.95), Inches(3.0), Inches(0.3), font_size=17, color=TITLE_COLOR, bold=True)
    for x, label in [
        (7.52, "Raw DICOM"),
        (9.25, "HU clip +\nnormalize"),
        (10.98, "Synthetic\nLDCT"),
    ]:
        add_rect(slide, Inches(x), Inches(3.52), Inches(1.35), Inches(0.92), WHITE, line=ACCENT)
        add_textbox(slide, label, Inches(x), Inches(3.64), Inches(1.35), Inches(0.5), font_size=15, color=TITLE_COLOR, bold=True, align=PP_ALIGN.CENTER)
    for x in [8.9, 10.63]:
        line = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, Inches(x), Inches(3.98), Inches(x + 0.28), Inches(3.98))
        line.line.color.rgb = ACCENT
        line.line.width = Pt(2.0)
        line.line.end_arrowhead = True
    add_textbox(slide, "Goal: credible, leakage-free evidence before claiming denoising gains.", Inches(7.52), Inches(4.85), Inches(4.45), Inches(0.55), font_size=15, color=BODY_COLOR)
    add_slide_number(slide, 2)

    # Slide 3
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Method: Diffusion as a Conservative Refiner")
    inputs = [
        (Inches(0.75), Inches(1.7), "Center LDCT"),
        (Inches(0.75), Inches(3.0), "Neighbor\nslices"),
        (Inches(0.75), Inches(4.3), "Edge map"),
    ]
    for x, y, label in inputs:
        add_rect(slide, x, y, Inches(1.65), Inches(0.8), GREEN_SOFT, line=RGBColor(160, 195, 170))
        add_textbox(slide, label, x, y + Inches(0.12), Inches(1.65), Inches(0.45), font_size=16, color=TITLE_COLOR, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(2.95), Inches(1.7), Inches(1.9), Inches(0.8), ACCENT_SOFT, line=ACCENT)
    add_textbox(slide, "RED-CNN\nanchor", Inches(2.95), Inches(1.84), Inches(1.9), Inches(0.45), font_size=18, color=TITLE_COLOR, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(3.15), Inches(2.95), Inches(2.8), Inches(1.2), WHITE, line=ACCENT)
    add_textbox(slide, "Edge-guided latent diffusion\nLDCT latent + neighbor context + Canny edge prior", Inches(3.22), Inches(3.12), Inches(2.66), Inches(0.78), font_size=16, color=TITLE_COLOR, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(6.55), Inches(2.95), Inches(2.25), Inches(1.2), RED_SOFT, line=GOLD)
    add_textbox(slide, "Anchor-aware blend\n+ edge-adaptive fusion", Inches(6.63), Inches(3.13), Inches(2.09), Inches(0.68), font_size=16, color=TITLE_COLOR, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(9.25), Inches(3.15), Inches(1.65), Inches(0.8), GREEN_SOFT, line=RGBColor(160, 195, 170))
    add_textbox(slide, "Output\nx-hat", Inches(9.25), Inches(3.27), Inches(1.65), Inches(0.45), font_size=17, color=TITLE_COLOR, bold=True, align=PP_ALIGN.CENTER)
    arrow_specs = [
        (2.4, 2.1, 2.95, 2.1),
        (2.4, 3.35, 3.15, 3.35),
        (2.4, 4.7, 3.15, 4.7),
        (4.85, 2.1, 6.55, 3.35),
        (5.95, 3.55, 6.55, 3.55),
        (8.8, 3.55, 9.25, 3.55),
    ]
    for x1, y1, x2, y2 in arrow_specs:
        line = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        line.line.color.rgb = ACCENT
        line.line.width = Pt(2.2)
        line.line.end_arrowhead = True
    add_bullets(
        slide,
        [
            "Key idea: the diffusion branch predicts a small structured residual, not the entire reconstruction.",
            "Edge guidance and 2.5D neighboring slices improve structural awareness.",
            "Auxiliary intensity and gradient losses translate structural gains into pixel fidelity.",
        ],
        Inches(0.95),
        Inches(5.45),
        Inches(11.4),
        Inches(1.45),
        font_size=17,
    )
    add_slide_number(slide, 3)

    # Slide 4
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Quantitative Results on the Held-Out Test Split")
    add_image_contain(slide, PAPER_FIGURES / "performance_overview.png", Inches(0.72), Inches(1.38), Inches(6.95), Inches(5.35))
    rows = [
        {"name": "LDCT", "psnr": "28.218", "ssim": "0.6244", "rmse": "0.08176", "highlight": False},
        {"name": "RED-CNN", "psnr": "35.020", "ssim": "0.9028", "rmse": "0.03598", "highlight": False},
        {"name": "No-Edge EG-LDM", "psnr": "27.681", "ssim": "0.7063", "rmse": "0.08328", "highlight": False},
        {"name": "Edge-Guided EG-LDM", "psnr": "28.463", "ssim": "0.7102", "rmse": "0.07652", "highlight": False},
        {"name": "Hybrid EG-LDM", "psnr": "35.120", "ssim": "0.9040", "rmse": "0.03558", "highlight": True},
    ]
    add_table_like_results(slide, rows, Inches(7.95), Inches(1.55), Inches(4.6), Inches(3.45))
    callout = add_rect(slide, Inches(7.95), Inches(5.3), Inches(4.62), Inches(1.05), ACCENT_SOFT, line=ACCENT)
    callout.shadow.inherit = False
    add_textbox(slide, "Hybrid vs RED-CNN\n+0.0998 dB PSNR | +0.00123 SSIM | -0.00040 RMSE", Inches(8.15), Inches(5.48), Inches(4.22), Inches(0.62), font_size=18, color=TITLE_COLOR, bold=True, align=PP_ALIGN.CENTER)
    add_slide_number(slide, 4)

    # Slide 5
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Qualitative Result and Final Takeaway")
    add_image_contain(slide, PAPER_FIGURES / "hybrid_comparison.png", Inches(0.7), Inches(1.15), Inches(12.0), Inches(2.15))
    add_textbox(slide, "Takeaways", Inches(0.85), Inches(3.55), Inches(2.5), Inches(0.35), font_size=18, color=ACCENT, bold=True)
    add_bullets(
        slide,
        [
            "Edge guidance improves the diffusion family over the no-edge ablation.",
            "Standalone diffusion remains soft; the best role of diffusion is hybrid refinement.",
            "Conservative anchor-aware and edge-adaptive fusion keeps the output tethered to measured LDCT.",
        ],
        Inches(0.9),
        Inches(3.92),
        Inches(7.2),
        Inches(2.1),
        font_size=18,
    )
    add_rect(slide, Inches(8.55), Inches(3.68), Inches(3.9), Inches(1.7), LIGHT, line=RGBColor(218, 224, 232))
    add_textbox(slide, "Limitation", Inches(8.78), Inches(3.9), Inches(2.0), Inches(0.3), font_size=18, color=ACCENT, bold=True)
    add_bullets(
        slide,
        [
            "Low-dose measurements are synthetic, not paired clinical acquisitions.",
            "Result gain over RED-CNN is modest but consistent under the reported protocol.",
        ],
        Inches(8.8),
        Inches(4.2),
        Inches(3.3),
        Inches(1.0),
        font_size=15,
    )
    add_textbox(slide, "Thank you.", Inches(10.85), Inches(6.6), Inches(1.5), Inches(0.3), font_size=20, color=TITLE_COLOR, bold=True, align=PP_ALIGN.RIGHT)
    add_slide_number(slide, 5)

    prs.save(str(OUTPUT_PPTX))

    script_text = """3-minute talk script
Title: Hybrid Edge-Guided Diffusion Refinement for Leakage-Free LDCT Denoising
Estimated duration: about 2 minutes 20 seconds at 135 words per minute

Speaker split
- Zhexi Feng: Slides 1-2
- Bingrui Zhang: Slides 3-5

Slide 1 - Zhexi Feng
Good afternoon. We study low-dose CT denoising under a strict question: can diffusion improve image quality without hallucinating anatomy? Our answer is yes, but only when diffusion is used conservatively, as a refiner around a strong deterministic RED-CNN anchor.

Slide 2 - Zhexi Feng
To make that claim credible, the protocol matters. We start from raw LIDC-IDRI DICOM, curate 243,957 valid CT slices from 1,010 patients, and enforce patient-level train, validation, and test splits. Because LIDC does not provide paired standard-dose and low-dose scans, we synthesize low-dose observations with a signal-dependent Gaussian corruption model. This gives us paired supervision without slice leakage.

Transition - Zhexi Feng
I will now hand it over to Bingrui for the model and the results.

Slide 3 - Bingrui Zhang
Our final model is hybrid. RED-CNN first produces a high-fidelity anchor. An edge-guided latent diffusion module then sees the noisy slice, neighboring slices, and a Canny edge map. At inference, anchor-aware blending and edge-adaptive fusion keep the final output close to the measured LDCT, so diffusion contributes only a small structured correction.

Slide 4 - Bingrui Zhang
The main quantitative result is on the held-out test split. Hybrid EG-LDM reaches 35.12 PSNR, 0.904 SSIM, and 0.0356 RMSE. RED-CNN reaches 35.02, 0.9028, and 0.0360. So the gain over the anchor is modest, but consistent. The standalone diffusion variants are much weaker, which tells us diffusion is not the best primary denoiser here.

Slide 5 - Bingrui Zhang
Qualitatively, the hybrid model preserves sharp boundaries while reducing local residual error. So our main takeaway is: edge guidance helps diffusion, but the strongest use of diffusion in this benchmark is conservative hybrid refinement. The main limitation is that the low-dose data are synthetic rather than truly paired clinical acquisitions. Thank you.
"""
    OUTPUT_SCRIPT.write_text(script_text)


if __name__ == "__main__":
    build_presentation()
